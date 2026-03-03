# -*- coding: utf-8 -*-
"""Build Fedi UX Audit .pptx -- fully editable in Google Slides."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

BLACK    = RGBColor(0x1A, 0x1A, 0x2E)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
RED      = RGBColor(0xE0, 0x3B, 0x3B)
ORANGE   = RGBColor(0xF0, 0x8C, 0x00)
TEAL     = RGBColor(0x00, 0x8B, 0x8B)
GREY     = RGBColor(0xF4, 0xF4, 0xF6)
DARKGREY = RGBColor(0x55, 0x55, 0x66)
ACCENT   = RGBColor(0x2D, 0x6A, 0xD6)

SCREENSHOTS = '/Users/annie/fedi-ux-audit-repo/screenshots/annotated/'
CRIT = chr(0x25CF)
MOD  = chr(0x25B3)
GOOD = chr(0x2713)

def new_prs():
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    return prs

def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def bg(slide, color=WHITE):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def box(slide, left, top, width, height, fill=None, line=None, lw=None):
    s = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
    if fill: s.fill.solid(); s.fill.fore_color.rgb = fill
    else: s.fill.background()
    if line: s.line.color.rgb = line
    if lw: s.line.width = Pt(lw)
    if not line: s.line.fill.background()
    return s

def txbox(slide, left, top, width, height,
          text='', bold=False, italic=False, size=12,
          color=BLACK, align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]; p.alignment = align
    run = p.add_run(); run.text = text
    run.font.bold = bold; run.font.italic = italic
    run.font.size = Pt(size); run.font.color.rgb = color
    return tb

def pill(slide, left, top, width, height, text, fill, text_color=WHITE, size=9):
    r = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
    r.fill.solid(); r.fill.fore_color.rgb = fill; r.line.fill.background()
    tf = r.text_frame; tf.word_wrap = False
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    run = p.add_run(); run.text = text
    run.font.bold = True; run.font.size = Pt(size); run.font.color.rgb = text_color
    return r

def header_bar(slide, title, subtitle=''):
    box(slide, 0, 0, 13.33, 1.0, fill=BLACK)
    txbox(slide, 0.3, 0.1, 9, 0.5, text=title, bold=True, size=20, color=WHITE)
    if subtitle:
        txbox(slide, 0.3, 0.58, 12.5, 0.35, text=subtitle,
              size=11, color=RGBColor(0xBB, 0xBB, 0xCC), italic=True)

def section_divider(prs, title, subtitle=''):
    slide = blank_slide(prs)
    bg(slide, BLACK)
    box(slide, 0, 0, 0.08, 7.5, fill=ACCENT)
    txbox(slide, 0.35, 2.95, 12, 0.9, text=title, bold=True, size=36, color=WHITE)
    if subtitle:
        txbox(slide, 0.35, 3.85, 11, 0.5, text=subtitle,
              size=14, color=RGBColor(0xAA, 0xBB, 0xCC))

def add_screenshot(slide, path, left, top, width, height):
    if os.path.exists(path):
        slide.shapes.add_picture(path, Inches(left), Inches(top), Inches(width), Inches(height))
    else:
        box(slide, left, top, width, height,
            fill=RGBColor(0xDD, 0xDD, 0xEE), line=RGBColor(0xAA, 0xAA, 0xBB), lw=0.5)

def finding_row(slide, left, y, width, severity, text, size=9.5):
    sev_c = {CRIT: RED, MOD: ORANGE, GOOD: TEAL}
    pill(slide, left, y, 0.18, 0.22, severity, sev_c.get(severity, DARKGREY), WHITE, 8)
    txbox(slide, left + 0.23, y, width - 0.23, 0.28, text=text, size=size, color=BLACK)

def add_table(slide, rows, cols, left, top, width, height):
    t = slide.shapes.add_table(rows, cols, Inches(left), Inches(top), Inches(width), Inches(height))
    return t.table

def set_cell(tbl, row, col, text, bold=False, size=8, color=BLACK,
             bg_color=None, align=PP_ALIGN.LEFT, italic=False):
    cell = tbl.cell(row, col)
    cell.text_frame.word_wrap = True
    if bg_color: cell.fill.solid(); cell.fill.fore_color.rgb = bg_color
    p = cell.text_frame.paragraphs[0]; p.alignment = align
    run = p.add_run(); run.text = text
    run.font.bold = bold; run.font.italic = italic
    run.font.size = Pt(size); run.font.color.rgb = color


# ── slides ────────────────────────────────────────────────────────────────────

def slide_title(prs):
    slide = blank_slide(prs)
    bg(slide, BLACK)
    box(slide, 0, 5.5, 13.33, 2.0, fill=ACCENT)
    box(slide, 0, 0, 0.12, 7.5, fill=ACCENT)
    txbox(slide, 0.5, 1.5, 12, 1.2, text='Fedi App', bold=True, size=52, color=WHITE)
    txbox(slide, 0.5, 2.7, 10, 0.7, text='UX Audit -- February 2026',
          bold=False, size=22, color=RGBColor(0x88, 0xAA, 0xFF))
    txbox(slide, 0.5, 3.5, 10, 0.5,
          text='Version 26.2.4  |  Android  |  Nielsen Norman Group heuristics',
          size=12, color=RGBColor(0xBB, 0xBB, 0xCC))

def slide_exec_summary(prs):
    slide = blank_slide(prs)
    bg(slide, WHITE)
    header_bar(slide, 'Executive Summary')
    box(slide, 0.4, 1.2, 12.5, 5.8, fill=GREY)
    box(slide, 0.4, 1.2, 0.07, 5.8, fill=ACCENT)
    tb = slide.shapes.add_textbox(Inches(0.65), Inches(1.35), Inches(12.0), Inches(5.5))
    tf = tb.text_frame; tf.word_wrap = True
    paras = [
        ("Fedi is trying to do something genuinely ambitious -- combine chat, payments, and mini apps inside a federated community model. That\'s a hard brief, and you can feel the vision behind it.", False, False),
        ('', False, False),
        ('But right now, the app is designed for people who already understand how it works.', True, False),
        ('', False, False),
        ("The onboarding drops users straight into federation jargon without explanation. Navigation is almost entirely unlabelled. The wallet -- Fedi's core feature -- is locked behind a joining process that can leave someone completely stranded.", False, False),
        ('', False, False),
        ('The visual design is clean and the intent is clear. But the gap between what the team understands and what a new user experiences is significant.', False, False),
        ('', False, False),
        ('These findings are organised in the order a real user would encounter them.', False, True),
    ]
    first = True
    for text, bold, italic in paras:
        if first: p = tf.paragraphs[0]; first = False
        else: p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run(); run.text = text
        run.font.bold = bold; run.font.italic = italic
        run.font.size = Pt(12.5); run.font.color.rgb = BLACK

def finding_slide(prs, num, title, subtitle, screenshot_path,
                  criticals, moderates, positives):
    slide = blank_slide(prs)
    bg(slide, WHITE)
    header_bar(slide, f'{num}. {title}', subtitle)
    add_screenshot(slide, screenshot_path, 0.3, 1.1, 4.2, 5.9)
    box(slide, 4.7, 1.1, 8.33, 5.9, fill=GREY)
    box(slide, 4.7, 1.1, 0.07, 5.9, fill=ACCENT)
    txbox(slide, 4.95, 1.2, 8.0, 0.3, text='Findings', bold=True, size=12, color=ACCENT)
    y = 1.6
    for t in criticals: finding_row(slide, 4.95, y, 8.0, CRIT, t); y += 0.4
    y += 0.05
    for t in moderates: finding_row(slide, 4.95, y, 8.0, MOD,  t); y += 0.4
    y += 0.05
    for t in positives: finding_row(slide, 4.95, y, 8.0, GOOD, t); y += 0.4

def slide_nav(prs):
    slide = blank_slide(prs)
    bg(slide, WHITE)
    header_bar(slide, '5. Core Navigation',
               'Four unlabelled bottom-nav icons. Every tab blocked by a promo sheet on open.')
    add_screenshot(slide, SCREENSHOTS + '06-home.png',       0.3,  1.1, 3.8, 5.9)
    add_screenshot(slide, SCREENSHOTS + '07-chat.png',       4.25, 1.1, 2.6, 2.8)
    txbox(slide, 4.25, 3.95, 2.6, 0.25, text='Chat', bold=True, size=9, color=DARKGREY)
    add_screenshot(slide, SCREENSHOTS + '08-mini-apps.png',  7.1,  1.1, 2.6, 2.8)
    txbox(slide, 7.1, 3.95, 2.6, 0.25, text='Mini Apps', bold=True, size=9, color=DARKGREY)
    box(slide, 9.9, 1.1, 3.2, 5.9, fill=GREY)
    box(slide, 9.9, 1.1, 0.06, 5.9, fill=ACCENT)
    findings = [
        (CRIT, "Four unlabelled icons -- users must guess each tab's purpose."),
        (CRIT, 'Promo bottom sheet blocks content on every tab switch.'),
        (CRIT, 'Wallet tab redirects to federation wall with no explanation.'),
        (MOD,  'Unlabelled QR icon on Fedi Global banner -- not obviously tappable.'),
        (MOD,  'Four unlabelled header icons appear on every screen.'),
        (MOD,  'Chat empty state is bare text -- no illustration or next step.'),
        (MOD,  'Mini apps icon grid: wildly inconsistent icon styles.'),
    ]
    sev_c = {CRIT: RED, MOD: ORANGE, GOOD: TEAL}
    y = 1.2
    for sev, text in findings:
        pill(slide, 9.95, y, 0.18, 0.22, sev, sev_c[sev], WHITE, 8)
        txbox(slide, 10.2, y, 2.85, 0.28, text=text, size=8.5, color=BLACK)
        y += 0.38

def slide_settings(prs):
    slide = blank_slide(prs)
    bg(slide, WHITE)
    header_bar(slide, '6. Settings & Profile',
               'A QR code with no label sits above everything else.')
    add_screenshot(slide, SCREENSHOTS + '09-settings-top.png',    0.3,  1.1, 3.1, 5.9)
    add_screenshot(slide, SCREENSHOTS + '10-settings-bottom.png', 3.55, 1.1, 2.5, 2.85)
    add_screenshot(slide, SCREENSHOTS + '11-settings-footer.png', 3.55, 4.15, 2.5, 2.85)
    box(slide, 6.2, 1.1, 6.9, 5.9, fill=GREY)
    box(slide, 6.2, 1.1, 0.06, 5.9, fill=ACCENT)
    findings = [
        (CRIT, 'QR code at top has no label -- wallet address? Login code? No way to know.'),
        (MOD,  '"#267f" hash suffix on display name -- technical noise for users.'),
        (MOD,  '"Nostr Details" listed with no explanation of what Nostr is.'),
        (MOD,  'External link icons used inconsistently instead of chevrons.'),
        (MOD,  '"Share logs" looks like plain text -- no affordance it is tappable.'),
        (MOD,  '"Fedimint 0.9.1" -- internal version noise for regular users.'),
        (GOOD, 'Settings are well-structured: Backup, PIN, Language, Currency.'),
        (GOOD, '"Ask Fedi" help entry is present in settings.'),
        (GOOD, 'Icon + label + chevron list pattern is clear and scannable.'),
    ]
    sev_c = {CRIT: RED, MOD: ORANGE, GOOD: TEAL}
    y = 1.2
    for sev, text in findings:
        pill(slide, 6.3, y, 0.18, 0.22, sev, sev_c[sev], WHITE, 8)
        txbox(slide, 6.55, y, 6.45, 0.28, text=text, size=9, color=BLACK)
        y += 0.4

def slide_priority_matrix(prs):
    slide = blank_slide(prs)
    bg(slide, WHITE)
    header_bar(slide, 'Priority Matrix', 'All findings ranked by heuristic and severity')
    rows_data = [
        ('Federation Discovery', 'No explanation of federation before asking users to join', 'H2', CRIT),
        ('Federation Discovery', 'Federation cards not tappable -- no detail before committing', 'H5', CRIT),
        ('Joining Flow',  'Error banner persists across multiple screens', 'H1', CRIT),
        ('Joining Flow',  '"View public federations" loops back to same screen', 'H3', CRIT),
        ('Joining Flow',  '"Maybe Later" disappears after first navigation', 'H3', CRIT),
        ('Joining Flow',  'Back button exits app -- no route to main UI if joining fails', 'H3', CRIT),
        ('Arriving in App', 'Display name auto-assigned without asking', 'H2', CRIT),
        ('Wallet Tab',    'Redirects to federation wall with no explanation', 'H1', CRIT),
        ('Core Nav',      'All four bottom nav tabs are unlabelled', 'H4', CRIT),
        ('Core Nav',      'Every tab triggers promo bottom sheet on first open', 'H8', CRIT),
        ('Settings',      'QR code at top of settings has no label', 'H6', CRIT),
        ('Joining Flow',  'Camera permission screen ambiguous (Continue = grant or defer?)', 'H5', MOD),
        ('Joining Flow',  'QR scanner has no instruction text', 'H6', MOD),
        ('Joining Flow',  'Error banner has no retry action', 'H9', MOD),
        ('Federation Discovery', 'Description text truncated mid-sentence', 'H8', MOD),
        ('Federation Discovery', 'Three tabs have no active state or explanation', 'H4', MOD),
        ('Arriving in App', '"#267f" hash suffix shown with no explanation', 'H2', MOD),
        ('Arriving in App', 'Home screen gives equal weight to three content types', 'H8', MOD),
        ('Chat Screen',   'Empty state is bare text -- no illustration or CTA', 'H7', MOD),
        ('Mini Apps',     'App icon grid has inconsistent icon styles', 'H4', MOD),
        ('Settings',      '"Nostr Details" unexplained', 'H2', MOD),
        ('Settings',      'External link icons vs chevrons inconsistent', 'H4', MOD),
        ('Splash Screen', 'Legal consent text is low contrast', 'H5', MOD),
        ('Splash Screen', 'Two-button layout is clear and familiar', 'H4', GOOD),
        ('Camera Permission', '"This can be updated later" reduces anxiety', 'H5', GOOD),
        ('Settings',      'Good range of settings options', 'H7', GOOD),
        ('Settings',      '"Ask Fedi" help entry present', 'H10', GOOD),
    ]
    n = len(rows_data) + 1
    tbl = add_table(slide, n, 4, 0.3, 1.1, 12.73, 6.1)
    tbl.columns[0].width = Inches(2.2)
    tbl.columns[1].width = Inches(8.5)
    tbl.columns[2].width = Inches(0.6)
    tbl.columns[3].width = Inches(0.75)
    for c, h in enumerate(['Area', 'Finding', 'H#', 'Sev.']):
        set_cell(tbl, 0, c, h, bold=True, size=9, color=WHITE,
                 bg_color=BLACK, align=PP_ALIGN.CENTER)
    sev_bg  = {CRIT: RGBColor(0xFF, 0xEB, 0xEB), MOD: RGBColor(0xFF, 0xF3, 0xCC), GOOD: RGBColor(0xE6, 0xF4, 0xF1)}
    sev_col = {CRIT: RED, MOD: ORANGE, GOOD: TEAL}
    for r, (area, finding, heur, sev) in enumerate(rows_data, 1):
        rb = GREY if r % 2 == 0 else WHITE
        set_cell(tbl, r, 0, area, size=8, color=BLACK, bg_color=rb)
        set_cell(tbl, r, 1, finding, size=8, color=BLACK, bg_color=rb)
        set_cell(tbl, r, 2, heur, size=8, color=DARKGREY, bg_color=rb, align=PP_ALIGN.CENTER)
        set_cell(tbl, r, 3, sev, bold=True, size=11, color=sev_col[sev],
                 bg_color=sev_bg[sev], align=PP_ALIGN.CENTER)

def slide_legend(prs):
    slide = blank_slide(prs)
    bg(slide, WHITE)
    header_bar(slide, 'Reference: Severity & Heuristics')
    box(slide, 0.4, 1.2, 4.5, 2.2, fill=GREY)
    txbox(slide, 0.6, 1.25, 4.0, 0.35, text='Severity Legend', bold=True, size=13, color=ACCENT)
    for y, (sym, col, label) in enumerate(
            [(CRIT, RED, 'Critical -- significant impact, prioritise immediately'),
             (MOD, ORANGE, 'Moderate -- noticeable friction, worth addressing'),
             (GOOD, TEAL, 'Works well -- effective design to preserve')], 0):
        yy = 1.65 + y * 0.42
        pill(slide, 0.55, yy, 0.28, 0.28, sym, col, WHITE, 10)
        txbox(slide, 0.9, yy, 3.8, 0.28, text=label, size=10, color=BLACK)
    box(slide, 5.2, 1.2, 7.9, 5.9, fill=GREY)
    txbox(slide, 5.4, 1.25, 7.5, 0.35,
          text='Nielsen Norman Group -- 10 Usability Heuristics', bold=True, size=13, color=ACCENT)
    heuristics = [
        ('H1', 'Visibility of system status'),
        ('H2', 'Match between system and the real world'),
        ('H3', 'User control and freedom'),
        ('H4', 'Consistency and standards'),
        ('H5', 'Error prevention'),
        ('H6', 'Recognition rather than recall'),
        ('H7', 'Flexibility and efficiency of use'),
        ('H8', 'Aesthetic and minimalist design'),
        ('H9', 'Help users recognise, diagnose, and recover from errors'),
        ('H10', 'Help and documentation'),
    ]
    y = 1.7
    for code, desc in heuristics:
        pill(slide, 5.35, y, 0.45, 0.26, code, ACCENT, WHITE, 8)
        txbox(slide, 5.88, y, 7.0, 0.26, text=desc, size=10, color=BLACK)
        y += 0.37



# ── Assumptions data ──────────────────────────────────────────────────────────

from pptx.dml.color import RGBColor as _C

ASSUMPTIONS = [
    dict(id='A1', theme='Federation vs Community mental model',
         assumption='Users understand the difference between Federation (money/bank) and Community (social/town square) quickly enough to use it correctly.',
         rqs=['What do users think each concept is?',
              'Where do they misclassify actions?',
              'What minimum explanation makes it click?',
              'Which labels or metaphors reduce confusion fastest?',
              'What errors happen when they misunderstand (lost funds, wrong help channel)?'],
         color=_C(0xD0, 0xE8, 0xFF)),
    dict(id='A2', theme='Onboarding and time-to-value',
         assumption='New users can reach their first successful moment (setup > join > receive/send) without external help.',
         rqs=['What is the first confusing step?',
              'What prevents first deposit or receive?',
              'What is the median time-to-first-value?',
              'What triggers dropout?',
              'What single piece of reassurance increases completion most?'],
         color=_C(0xE6, 0xF4, 0xE6)),
    dict(id='A3', theme='Social Backup trust',
         assumption='Social Backup feels safer (or at least acceptable) vs self-custody for non-technical users.',
         rqs=['What do users think Social Backup does?',
              'What are their nightmare scenarios (theft, coercion, family conflict)?',
              'Which social structures make it viable (elders, church, co-op, family)?',
              'What proof or trust cues do they need before enabling it?'],
         color=_C(0xFF, 0xE0, 0xE0)),
    dict(id='A4', theme='Recovery and panic moments',
         assumption='In loss or panic scenarios (lost phone, reinstall), users can recover confidently without catastrophic errors.',
         rqs=['What is the user mental model in panic?',
              'Which recovery steps are misinterpreted?',
              'What reassurance reduces churn and support contacts?',
              'Which safeguards prevent irreversible actions?'],
         color=_C(0xFF, 0xE0, 0xE0)),
    dict(id='A5', theme='Offline transactions',
         assumption='Offline payments are both usable and trusted in patchy-connectivity contexts.',
         rqs=['When do users actually need offline payments?',
              'What failure modes are acceptable (delayed sync, pending state)?',
              'How do users resolve disputes over offline payments?',
              'What UI cues prevent false confidence?'],
         color=_C(0xE6, 0xF4, 0xE6)),
    dict(id='A6', theme='Stable Balance (volatility protection)',
         assumption='Users want a stable-value option and correctly understand the trade-off (stability vs bitcoin amount).',
         rqs=['Do users understand what stays constant vs what changes?',
              'When do they choose Stable vs Bitcoin balance?',
              'Does it increase day-to-day use or add cognitive load?',
              'Does it create two-money confusion?'],
         color=_C(0xE6, 0xF4, 0xE6)),
    dict(id='A7', theme='Mini apps value',
         assumption='Mini apps provide genuine daily utility rather than clutter or confusion.',
         rqs=['Which mini apps are must-have by region?',
              'How do users assess trustworthiness of apps?',
              'Does the catalog feel safe or sketchy?',
              'What is the smallest set of apps that drives weekly retention?'],
         color=_C(0xE6, 0xF4, 0xE6)),
    dict(id='A8', theme='In-app money and chat',
         assumption='Sending money in chat is intuitive and increases usage vs separate wallet actions.',
         rqs=['Do users understand who is being paid and from which balance?',
              'What mistakes happen (wrong person, wrong amount)?',
              'Does chat increase trust or increase scam risk?'],
         color=_C(0xD0, 0xE8, 0xFF)),
    dict(id='A9', theme='Trust and scam resistance',
         assumption='Users can avoid common scam patterns and understand safety boundaries without becoming paranoid.',
         rqs=['What are the top scam vectors in each region?',
              'Which UI friction is acceptable before it harms adoption?',
              'What social proof signals are credible?',
              'Which safety education moments stick?'],
         color=_C(0xFF, 0xE0, 0xE0)),
    dict(id='A10', theme='Language, literacy and comprehension',
         assumption='The product concepts can be understood by users with varying literacy and education levels.',
         rqs=['Which terms fail comprehension (federation, backup, guardian)?',
              'What metaphors work locally?',
              'Do icons help or confuse?',
              'What is the minimum viable explanation for each high-stakes concept?'],
         color=_C(0xD0, 0xE8, 0xFF)),
    dict(id='A11', theme='Device sharing and privacy',
         assumption='The app remains safe and usable in device-sharing households (common in some contexts).',
         rqs=['Do users share phones?',
              'What content should be hidden by default?',
              'What privacy modes are needed?',
              'What threats exist from family or community power dynamics?'],
         color=_C(0xFF, 0xE0, 0xE0)),
    dict(id='A12', theme='Community governance and support',
         assumption='Communities and federations can provide support and governance without causing confusion about accountability.',
         rqs=['When something goes wrong, who do users believe is responsible?',
              'What support path do they naturally take?',
              'What UX makes where-to-get-help unambiguous?'],
         color=_C(0xFF, 0xF3, 0xCC)),
    dict(id='A13', theme='Cross-region generalisation',
         assumption='Some UX patterns generalise across regions; others must be localised.',
         rqs=['Which behaviours are universal (trust cues, fear triggers)?',
              'Which are regional (financial rails, norms)?',
              'What is the minimum sample strategy to identify universal vs local?',
              'How do we avoid overfitting to one market?'],
         color=_C(0xFF, 0xF3, 0xCC)),
    dict(id='A14', theme='Western vs Global South segmentation',
         assumption='Western users and Global South users need different onboarding, terminology, or defaults -- but not two totally different products.',
         rqs=['Which steps differ most by segment?',
              'What defaults should change (privacy, backup prompts, education)?',
              'Can a single flow adapt with progressive disclosure?',
              'Where do segments share core needs?'],
         color=_C(0xFF, 0xF3, 0xCC)),
    dict(id='A15', theme='Evidence threshold',
         assumption='The org can define validated-enough evidence thresholds per decision so shipping is not blocked by perfection.',
         rqs=['Which decisions require high confidence (recovery, custody)?',
              'Which can ship as experiments?',
              'What counts as success (task success, comprehension, trust rating)?',
              'How do we express confidence levels to engineers and leadership?'],
         color=_C(0xFF, 0xF3, 0xCC)),
    dict(id='A16', theme='Roadmap alignment',
         assumption='Research effort must align to Product and Engineering plans; some work is in-flight and needs harm-reduction rather than discovery.',
         rqs=['What is already committed vs shapeable?',
              'Where can research still reduce risk mid-flight?',
              'What research will unlock near-term roadmap decisions fastest?'],
         color=_C(0xFF, 0xF3, 0xCC)),
]


def slide_assumptions_overview(prs):
    slide = blank_slide(prs)
    bg(slide, WHITE)
    header_bar(slide, 'Assumptions + Research Questions',
               '16 testable claims -- each should be validated before committing to design decisions')
    txbox(slide, 0.4, 1.1, 12.5, 0.45,
          text=("These assumptions represent the team's current mental model of how users think and behave. "
                'Each one is a testable hypothesis -- not a given.'),
          size=11, color=DARKGREY, italic=True)
    legend = [
        ('Mental model', _C(0xD0, 0xE8, 0xFF)),
        ('Trust & safety', _C(0xFF, 0xE0, 0xE0)),
        ('Onboarding & flows', _C(0xE6, 0xF4, 0xE6)),
        ('Strategy & governance', _C(0xFF, 0xF3, 0xCC)),
    ]
    x = 0.4
    for label, col in legend:
        box(slide, x, 1.65, 0.18, 0.18, fill=col, line=_C(0xCC, 0xCC, 0xCC), lw=0.5)
        txbox(slide, x + 0.22, 1.63, 2.8, 0.22, text=label, size=9, color=DARKGREY)
        x += 3.15
    card_w, card_h, gap_x, gap_y = 3.0, 1.55, 0.17, 0.13
    for i, a in enumerate(ASSUMPTIONS):
        col = i % 4
        row = i // 4
        x = 0.32 + col * (card_w + gap_x)
        y = 2.0 + row * (card_h + gap_y)
        box(slide, x, y, card_w, card_h, fill=a['color'], line=_C(0xBB, 0xBB, 0xBB), lw=0.5)
        pill(slide, x + 0.06, y + 0.06, 0.32, 0.22, a['id'], BLACK, WHITE, 7)
        txbox(slide, x + 0.42, y + 0.06, card_w - 0.48, 0.22,
              text=a['theme'], bold=True, size=7.5, color=BLACK)
        txbox(slide, x + 0.08, y + 0.31, card_w - 0.14, 1.2,
              text=a['assumption'], size=7, color=DARKGREY, italic=True)


def slides_rq_detail(prs):
    chunks = [ASSUMPTIONS[i:i+4] for i in range(0, len(ASSUMPTIONS), 4)]
    for idx, chunk in enumerate(chunks):
        slide = blank_slide(prs)
        bg(slide, WHITE)
        lo = idx * 4 + 1
        hi = min(idx * 4 + 4, len(ASSUMPTIONS))
        header_bar(slide,
                   'Assumptions + Research Questions  (A%d--A%d)' % (lo, hi),
                   'Each assumption is a testable hypothesis. Questions show how to validate or falsify it.')
        card_w = (13.33 - 0.5) / len(chunk) - 0.15
        for ci, a in enumerate(chunk):
            x = 0.25 + ci * (card_w + 0.15)
            y = 1.1
            box(slide, x, y, card_w, 6.2, fill=a['color'], line=_C(0xCC, 0xCC, 0xCC), lw=0.5)
            box(slide, x, y, card_w, 0.55, fill=BLACK)
            pill(slide, x + 0.08, y + 0.14, 0.38, 0.25, a['id'], ACCENT, WHITE, 8)
            txbox(slide, x + 0.5, y + 0.08, card_w - 0.55, 0.38,
                  text=a['theme'], bold=True, size=8, color=WHITE)
            txbox(slide, x + 0.1, y + 0.62, card_w - 0.2, 1.3,
                  text=a['assumption'], size=8, color=BLACK, italic=True)
            box(slide, x + 0.08, y + 2.0, card_w - 0.16, 0.25, fill=ACCENT)
            txbox(slide, x + 0.1, y + 2.0, card_w - 0.2, 0.25,
                  text='Research Questions', bold=True, size=7.5, color=WHITE)
            rq_y = y + 2.3
            for rq in a['rqs']:
                txbox(slide, x + 0.18, rq_y, card_w - 0.28, 0.48,
                      text='- ' + rq, size=7.5, color=BLACK)
                rq_y += 0.5


# ── Assumption Register ───────────────────────────────────────────────────────

REGISTER = [
    dict(id='A1', theme='Core mental model',
         assumption='Users understand Federation vs Community well enough to act correctly.',
         risk='H', evidence='None', region='Global',
         decision='IA, labeling, onboarding, navigation',
         method='Comprehension + task-based usability on prototype',
         sample='8-12/region; new users', timebox='5-7 days',
         notes='High-risk confusion = wrong actions + loss of trust'),
    dict(id='A2', theme='Onboarding',
         assumption='Users can reach first successful moment without external help.',
         risk='H', evidence='None', region='Global',
         decision='Onboarding steps, defaults, prompts',
         method='Funnel review + 6-10 usability sessions',
         sample='10-15 total; new users', timebox='1-2 weeks',
         notes='Pair qual with funnel data if instrumentation exists'),
    dict(id='A3', theme='Social Backup trust',
         assumption='Social Backup feels safe + socially acceptable to target users.',
         risk='H', evidence='None', region='Regional',
         decision='Recovery design, messaging, rollout strategy',
         method='Concept test + scenario probes',
         sample='6-8/region; mix gender/age', timebox='1-2 weeks',
         notes='Include power dynamics (family/community coercion)'),
    dict(id='A4', theme='Recovery / panic',
         assumption='In a panic moment, users can recover without catastrophic errors.',
         risk='H', evidence='None', region='Global',
         decision='Recovery flow, safeguards, copy, support',
         method='Disaster drill usability (lost phone scenario)',
         sample='10-12; low-tech users included', timebox='1-2 weeks',
         notes='One-way-door decision -- highest evidence bar'),
    dict(id='A5', theme='Offline transactions',
         assumption='Offline payments are usable and trusted in real contexts.',
         risk='H', evidence='None', region='Regional',
         decision='Offline UX, sync states, dispute handling',
         method='Field/diary + scenario usability',
         sample='6-10 in offline-relevant context', timebox='2-3 weeks',
         notes='Must test trust of pending/sync-later states'),
    dict(id='A6', theme='Stable Balance',
         assumption='Users want stable value and understand the trade-off.',
         risk='M/H', evidence='None', region='Regional',
         decision='Balance UX, education, default settings',
         method='Concept comprehension + choice tasks',
         sample='8-12; mix bitcoin-native + general', timebox='1-2 weeks',
         notes='High risk if it creates two-money confusion'),
    dict(id='A7', theme='Mini apps value',
         assumption='Mini apps drive weekly utility, not clutter.',
         risk='M', evidence='None', region='Regional',
         decision='Catalog IA, prioritization, governance',
         method='Jobs-to-be-done + card sort',
         sample='10-15; active community members', timebox='2 weeks',
         notes='Include trust-in-catalog vs random links'),
    dict(id='A8', theme='Money in chat',
         assumption='Paying in chat is clear + reduces friction.',
         risk='M', evidence='None', region='Global',
         decision='Chat-pay UX, confirmations, error prevention',
         method='Usability + mis-send prevention test',
         sample='10-12; novice users', timebox='1 week',
         notes='Prioritize wrong person/amount prevention'),
    dict(id='A9', theme='Scam resistance',
         assumption='Users can avoid common scams without UX becoming oppressive.',
         risk='H', evidence='Weak', region='Regional',
         decision='Warnings, confirmations, education, trust cues',
         method='Threat-model workshop + rapid user validation',
         sample='6-8 users + internal stakeholders', timebox='1-2 weeks',
         notes='Combine user input + support/community signals'),
    dict(id='A10', theme='Language & literacy',
         assumption='Key concepts can be understood across literacy levels.',
         risk='H', evidence='None', region='Regional',
         decision='Terminology, iconography, education',
         method='Comprehension testing w/ low-literacy variants',
         sample='8-10/region; varying literacy', timebox='2 weeks',
         notes='Include teach-back method'),
    dict(id='A11', theme='Privacy / device sharing',
         assumption='App remains safe in device-sharing environments.',
         risk='H', evidence='None', region='Regional',
         decision='Privacy defaults, lock modes, content visibility',
         method='Context interviews + privacy prototype test',
         sample='6-10 in relevant contexts', timebox='2-3 weeks',
         notes='Safety-critical; may vary hugely by locale'),
    dict(id='A12', theme='Support & accountability',
         assumption='Users know who to blame / where to get help.',
         risk='M/H', evidence='None', region='Global',
         decision='Support UX, escalation paths, messaging',
         method='Journey mapping + support flow usability',
         sample='10-12; mix new/active', timebox='1-2 weeks',
         notes='Reduces churn + support load; improves trust'),
    dict(id='A13', theme='Cross-region generalisation',
         assumption='Some patterns generalise; others need localisation.',
         risk='H', evidence='None', region='Global',
         decision='Research strategy, product strategy',
         method='Comparative synthesis across region cycles',
         sample='N/A (meta)', timebox='Ongoing',
         notes='This is the research-OS guardrail assumption'),
    dict(id='A14', theme='Evidence threshold',
         assumption='Org can define validated-enough thresholds per decision type.',
         risk='H', evidence='Weak', region='Global',
         decision='Shipping governance',
         method='Define validation ladder + align with PM/Eng',
         sample='Stakeholders', timebox='1 week',
         notes='Makes the prove-it request scalable'),
]


def _register_slide(prs, rows, suffix=''):
    slide = blank_slide(prs)
    bg(slide, WHITE)
    header_bar(slide, 'Assumption Register' + suffix,
               'Prioritised research backlog: risk, evidence status, method & timebox')
    n = len(rows) + 1
    tbl = add_table(slide, n, 9, 0.15, 1.1, 13.0, 6.15)
    cols   = ['ID', 'Theme', 'Assumption', 'Risk', 'Evidence', 'Region', 'Decision', 'Method', 'Timebox']
    widths = [0.35, 1.1, 2.6, 0.38, 0.7, 0.65, 1.45, 2.6, 0.7]
    for i, w in enumerate(widths):
        tbl.columns[i].width = Inches(w)
    for c, h in enumerate(cols):
        set_cell(tbl, 0, c, h, bold=True, size=8, color=WHITE,
                 bg_color=BLACK, align=PP_ALIGN.CENTER)
    risk_bg  = {'H': _C(0xFF,0xEB,0xEB), 'M': _C(0xFF,0xF3,0xCC),
                'L': _C(0xE6,0xF4,0xE6), 'M/H': _C(0xFF,0xEC,0xD5)}
    risk_col = {'H': RED, 'M': ORANGE, 'L': TEAL, 'M/H': _C(0xD0,0x60,0x00)}
    ev_bg    = {'None': _C(0xFF,0xEB,0xEB), 'Weak': _C(0xFF,0xF3,0xCC),
                'Mod':  _C(0xE8,0xF0,0xFE), 'Strong': _C(0xE6,0xF4,0xE6)}
    for r, row in enumerate(rows, 1):
        rb   = GREY if r % 2 == 0 else WHITE
        risk = row['risk']
        ev   = row['evidence']
        set_cell(tbl, r, 0, row['id'],        bold=True, size=8, color=ACCENT, bg_color=rb, align=PP_ALIGN.CENTER)
        set_cell(tbl, r, 1, row['theme'],      bold=True, size=7.5, color=BLACK, bg_color=rb)
        set_cell(tbl, r, 2, row['assumption'], size=7, color=BLACK, bg_color=rb)
        set_cell(tbl, r, 3, risk, bold=True, size=9, color=risk_col.get(risk, BLACK),
                 bg_color=risk_bg.get(risk, rb), align=PP_ALIGN.CENTER)
        set_cell(tbl, r, 4, ev, size=7.5, color=DARKGREY,
                 bg_color=ev_bg.get(ev, rb), align=PP_ALIGN.CENTER)
        set_cell(tbl, r, 5, row['region'],   size=7.5, color=DARKGREY, bg_color=rb, align=PP_ALIGN.CENTER)
        set_cell(tbl, r, 6, row['decision'], size=7, color=BLACK, bg_color=rb)
        set_cell(tbl, r, 7, row['method'],   size=7, color=BLACK, bg_color=rb)
        set_cell(tbl, r, 8, row['timebox'],  size=7.5, color=DARKGREY, bg_color=rb, align=PP_ALIGN.CENTER)


def slides_assumption_register(prs):
    _register_slide(prs, REGISTER[:8],  ' (A1-A8)')
    _register_slide(prs, REGISTER[8:],  ' (A9-A14)')


# ── main ──────────────────────────────────────────────────────────────────────



# ── Speaker notes helper ──────────────────────────────────────────────────────

def set_notes(slide, text):
    """Set speaker notes on a slide."""
    tf = slide.notes_slide.notes_text_frame
    tf.text = text


# ── 30/60/90 Days section ─────────────────────────────────────────────────────

def slide_30(prs):
    slide = blank_slide(prs)
    bg(slide, WHITE)
    header_bar(slide, 'First 30 Days: Build credibility + install the operating system',
               'Goal: be trusted, be organised, and deliver one fast win')

    # 4 output cards
    cards = [
        ('Roadmap-to-Research Map (v0)',
         'One page mapping initiatives → assumptions → key decisions → proposed study types.',
         ACCENT),
        ('Assumption Register + RQ Bank (v0)',
         'Lightweight but real. The "where research comes from" source of truth.',
         RGBColor(0x2D, 0x9A, 0x6A)),
        ('1-2 Fast Tactical Studies',
         'Pick something currently painful -- onboarding, recovery, mini-app discovery -- and ship 3-5 actionable fixes fast.',
         ORANGE),
        ('Research Ops Starter Kit',
         'Templates: discussion guide, usability script, consent, note-taking, insight write-up.',
         RGBColor(0x7B, 0x2D, 0xBF)),
    ]
    card_w = 2.9
    for i, (title, body, col) in enumerate(cards):
        x = 0.35 + i * (card_w + 0.2)
        box(slide, x, 1.2, card_w, 5.9, fill=RGBColor(0xF7, 0xF8, 0xFF),
            line=RGBColor(0xCC, 0xCC, 0xDD), lw=0.5)
        box(slide, x, 1.2, card_w, 0.08, fill=col)
        pill(slide, x + 0.12, 1.35, 0.45, 0.28, str(i + 1), col, WHITE, 11)
        txbox(slide, x + 0.65, 1.35, card_w - 0.75, 0.4,
              text=title, bold=True, size=10.5, color=BLACK)
        txbox(slide, x + 0.12, 1.85, card_w - 0.24, 5.1,
              text=body, size=10, color=DARKGREY, wrap=True)

    set_notes(slide,
        "30-Day framing for the hiring conversation:\n\n"
        "1. Roadmap-to-Research Map: before running any study, build a one-pager that connects Fedi's product "
        "initiatives to the 16 assumptions in the register and the decisions those assumptions inform. "
        "This becomes the negotiation tool with PM/Eng about priority.\n\n"
        "2. Assumption Register v0: the lightweight version is already done (this deck). "
        "The live version gets owned in Notion/Sheets and updated after every study.\n\n"
        "3. Fast tactical win: Anwar has already identified the onboarding wall and the federation mental model "
        "as the most painful. One 6-session usability study with a mid-fi prototype can ship 3-5 specific "
        "copy/flow fixes within 2 weeks of starting.\n\n"
        "4. Research Ops starter kit: Anwar has built this from scratch at previous orgs (0 to 50+ structured "
        "studies). At Fedi this means: a reusable discussion guide template, consent form, note-taking sheet, "
        "and a 1-page insight write-up format -- so studies can be run or observed by PMs without bottlenecking "
        "the researcher."
    )
    return slide


def slide_60(prs):
    slide = blank_slide(prs)
    bg(slide, WHITE)
    header_bar(slide, 'Days 31-60: Prove multi-region research without burning the calendar',
               'Goal: build comparative insight across regions; make recruitment self-sustaining')

    outputs = [
        ('Regional Thin-Slice Cycles',
         [
             'Remote-first, community-assisted format',
             '6-8 interviews per region archetype',
             '+ 1 evaluative prototype test where relevant',
             'Builds the universal vs context-specific map',
         ],
         RGBColor(0xD0, 0xE8, 0xFF)),
        ('Recruitment Flywheel',
         [
             'Partner with federation/community admins',
             'Continuous warm pipeline -- no cold starts',
             'Screener + incentive framework in place',
             'Target: 0 weeks recruitment lag by Day 60',
         ],
         RGBColor(0xE6, 0xF4, 0xE6)),
        ('Monthly Evidence Review (30 min)',
         [
             'Standing slot with Product + Eng leads',
             'Format: what we learned / what changed / what we still don\'t know',
             'Research shapes sprint priorities, not post-hoc',
             'Creates accountability without bureaucracy',
         ],
         RGBColor(0xFF, 0xF3, 0xCC)),
    ]

    card_w = 3.9
    for i, (title, bullets, col) in enumerate(outputs):
        x = 0.35 + i * (card_w + 0.27)
        box(slide, x, 1.2, card_w, 5.9, fill=col,
            line=RGBColor(0xCC, 0xCC, 0xCC), lw=0.5)
        txbox(slide, x + 0.15, 1.3, card_w - 0.3, 0.45,
              text=title, bold=True, size=12, color=BLACK)
        box(slide, x + 0.15, 1.78, card_w - 0.3, 0.03, fill=RGBColor(0xBB, 0xBB, 0xBB))
        y = 1.9
        for b in bullets:
            txbox(slide, x + 0.3, y, card_w - 0.45, 0.5,
                  text=chr(0x25B8) + '  ' + b, size=10, color=BLACK)
            y += 0.52

    set_notes(slide,
        "Days 31-60 framing:\n\n"
        "Regional thin-slice cycles: the key insight here is that Fedi's target markets (Latin America, "
        "Africa, South/Southeast Asia) share some UX failure modes but not all. A thin-slice approach -- "
        "6-8 interviews per region archetype rather than exhaustive sampling -- lets you build the "
        "universal/regional map within a quarter, not a year. Pair with one prototype test per region "
        "where a specific flow is in scope (e.g. onboarding, recovery). Remote-first with community-assisted "
        "recruitment means costs stay low and ecological validity stays high.\n\n"
        "Recruitment flywheel: the biggest time sink for a solo researcher is participant recruitment. "
        "Partnering with federation and community admins to maintain a warm pipeline solves this. "
        "Screener + incentive framework should be standardised by Day 45 so studies can spin up within days.\n\n"
        "Monthly evidence review: 30-minute standing slot with PM and Eng leads. "
        "Not a report -- a conversation. Format: 3 sections, what we learned this month, "
        "what we changed as a result, what we still don\'t know. "
        "This prevents research from becoming a post-hoc justification and turns it into a sprint input."
    )
    return slide


def slide_90(prs):
    slide = blank_slide(prs)
    bg(slide, WHITE)
    header_bar(slide, 'Days 61-90: Turn research into compounding advantage',
               'Goal: research shapes the roadmap, not the other way around')

    # Timeline bar
    box(slide, 0.4, 1.25, 12.3, 0.12, fill=ACCENT)
    for i, (label, x) in enumerate([('Day 1', 0.35), ('Day 30', 3.5), ('Day 60', 7.0), ('Day 90', 10.5)]):
        box(slide, x + 0.1, 1.17, 0.05, 0.28, fill=BLACK)
        txbox(slide, x, 1.48, 1.2, 0.25, text=label, size=8.5,
              color=ACCENT if i == 3 else DARKGREY, bold=(i == 3))

    outputs = [
        ('Quarterly Research Roadmap (v1)',
         'Aligned to Product/Eng sprint cycles. Explicit tradeoffs: what you\'re studying, what you\'re parking and why. Living doc, updated monthly.',
         ACCENT),
        ('Decision-Grade Benchmarks',
         'For highest-risk assumptions (A1, A2, A3, A4), define evidence bars: e.g. 80% task success, comprehension threshold, trust signal threshold. Makes "is this good enough?" answerable.',
         RED),
        ('Self-Serve Enablement',
         'Train PMs and designers to run safe studies (concept sorting, lightweight usability checks) while you focus on high-stakes strategic work. Multiplies research output without multiplying headcount.',
         RGBColor(0x2D, 0x9A, 0x6A)),
    ]

    card_w = 3.9
    for i, (title, body, col) in enumerate(outputs):
        x = 0.35 + i * (card_w + 0.27)
        box(slide, x, 1.9, card_w, 5.2, fill=GREY,
            line=RGBColor(0xCC, 0xCC, 0xCC), lw=0.5)
        box(slide, x, 1.9, card_w, 0.08, fill=col)
        txbox(slide, x + 0.15, 2.05, card_w - 0.3, 0.45,
              text=title, bold=True, size=11, color=BLACK)
        txbox(slide, x + 0.15, 2.6, card_w - 0.3, 4.3,
              text=body, size=10, color=DARKGREY, wrap=True)

    set_notes(slide,
        "Days 61-90 framing:\n\n"
        "Quarterly Research Roadmap: by Day 90 you should own a living research roadmap that sits "
        "alongside the product roadmap. It shows which assumptions are being tested this quarter, "
        "which are parked (and why), and what the expected outputs are. This is the artefact that "
        "makes research a strategic input rather than a reactive service.\n\n"
        "Decision-grade benchmarks: for the highest-risk assumptions (mental model confusion, onboarding "
        "completion, Social Backup trust, recovery panic), define explicit evidence bars before running "
        "the study. Examples: 80% task success on a recovery flow = validated; below 60% = redesign "
        "required. This makes the 'prove it' question answerable in advance, so shipping decisions "
        "aren\'t blocked by open-ended debate.\n\n"
        "Self-serve enablement: a solo researcher cannot cover everything. The leverage move is training "
        "PMs and designers to run 'safe' studies -- concept sorting, 5-second tests, lightweight "
        "usability checks -- while the researcher focuses on complex, high-stakes, multi-region work. "
        "This is the exact model Anwar has executed at previous orgs, scaling from 0 to 50+ structured "
        "studies while maintaining quality gates."
    )
    return slide


def slide_9090_summary(prs):
    """One-page 30/60/90 overview timeline."""
    slide = blank_slide(prs)
    bg(slide, WHITE)
    header_bar(slide, '30 / 60 / 90 Day Plan -- Overview',
               'Designed for solo researcher + fast shipping environment')

    phases = [
        ('30', 'Build credibility + install the OS',
         ['Roadmap-to-Research map', 'Assumption Register + RQ Bank', '1-2 fast tactical studies', 'Research Ops starter kit'],
         ACCENT),
        ('60', 'Prove multi-region research at pace',
         ['Regional thin-slice cycles (6-8/region)', 'Recruitment flywheel with community admins', 'Monthly 30-min evidence review with PM/Eng'],
         RGBColor(0x2D, 0x9A, 0x6A)),
        ('90', 'Turn research into compounding advantage',
         ['Quarterly Research Roadmap v1', 'Decision-grade evidence benchmarks', 'Self-serve enablement for PMs/designers'],
         RGBColor(0x7B, 0x2D, 0xBF)),
    ]

    # Connecting line
    box(slide, 1.55, 3.0, 9.8, 0.08, fill=RGBColor(0xDD, 0xDD, 0xDD))

    card_w = 3.6
    for i, (days, title, bullets, col) in enumerate(phases):
        x = 0.5 + i * (card_w + 0.6)
        # Day circle
        c = slide.shapes.add_shape(9, Inches(x + 1.1), Inches(2.6), Inches(1.3), Inches(1.3))
        c.fill.solid(); c.fill.fore_color.rgb = col; c.line.fill.background()
        tf = c.text_frame; tf.word_wrap = False
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        run = p.add_run(); run.text = 'Day ' + days
        run.font.bold = True; run.font.size = Pt(13); run.font.color.rgb = WHITE

        txbox(slide, x, 4.05, card_w, 0.5,
              text=title, bold=True, size=11, color=col, align=PP_ALIGN.CENTER)
        y = 4.65
        for b in bullets:
            txbox(slide, x + 0.15, y, card_w - 0.3, 0.45,
                  text=chr(0x25B8) + '  ' + b, size=9.5, color=BLACK)
            y += 0.48

    set_notes(slide,
        "This slide gives the hiring panel a single-page summary of the phased approach.\n\n"
        "The underlying logic: Fedi needs research to be both rigorous and fast. "
        "A solo researcher can achieve this by:\n"
        "- Front-loading the operating system (register, roadmap, ops kit) so each study costs less\n"
        "- Using thin-slice regional cycles to build comparative insight without exhaustive sampling\n"
        "- Defining evidence bars upfront so 'good enough' is answerable, not debatable\n"
        "- Enabling self-serve at the margin so the researcher focuses on the hard strategic work\n\n"
        "This is the same shape as Anwar's previous UXR practice builds -- phased, pragmatic, "
        "balancing rigour with real constraints."
    )
    return slide


# ── main ──────────────────────────────────────────────────────────────────────

def build():
    prs = new_prs()

    slide_title(prs)
    set_notes(list(prs.slides)[-1],
        'Fedi App UX Audit -- February 2026\nVersion 26.2.4 | Android emulator | Nielsen Norman Group heuristics\nAudit conducted via live emulator walkthrough.')

    slide_exec_summary(prs)
    set_notes(list(prs.slides)[-1],
        'The core tension: Fedi is a technically sophisticated product trying to serve users who are not '
        'technically sophisticated. The onboarding wall (federation jargon, unlabelled navigation, wallet '
        'locked behind joining) creates a first-experience failure that undermines everything downstream. '
        'The findings are ordered by user encounter, not severity, to make the narrative coherent for '
        'stakeholders who haven\'t used the app.')

    section_divider(prs, 'UX Audit Findings',
                    'Organised in the order a real user encounters them')

    finding_slide(
        prs, 1, 'First Launch', 'A calm, well-designed splash screen.',
        SCREENSHOTS + '01-splash.png',
        criticals=[],
        moderates=[
            'Legal copy ("By continuing, you agree...") is low contrast and easy to miss. For a financial app, this matters.',
            '"Ask Fedi" help link is so peripheral most users will not notice it.',
        ],
        positives=[
            'Clear visual hierarchy, intuitive two-button layout.',
            'Avatar bubbles communicate social-network context without words.',
        ],
    )
    set_notes(list(prs.slides)[-1],
        'First impressions are broadly positive. The splash screen does its job -- clean, calm, '
        'minimal cognitive load. The two issues are both fixable without design rework:\n'
        '- Legal copy: increase contrast, move below the CTA button (standard pattern)\n'
        '- Ask Fedi: consider surfacing as a persistent help icon (?) in the nav rather than a buried footer link\n\n'
        'Heuristic: H5 (error prevention), H10 (help and documentation)')

    finding_slide(
        prs, 2, 'Federation Discovery -- The First Wall',
        'They tap Get Started and land on Join or Create a Federation. No one has explained what a federation is.',
        SCREENSHOTS + '02-federation-discovery.png',
        criticals=[
            'The word federation appears repeatedly before the user has any idea what it means.',
            'Federation list cards are not tappable -- no detail view before committing to join.',
        ],
        moderates=[
            'Description text truncated mid-sentence ("Privacy Sovereignty Fr...").',
            'Three tabs (Discover, Join, Create) have no active state and no explanation.',
        ],
        positives=['"Maybe Later" at the bottom gives users a safe exit.'],
    )
    set_notes(list(prs.slides)[-1],
        'This is the highest-friction moment in the entire flow. The user has just downloaded a financial '
        'app and is immediately asked to join something called a federation -- with no explanation of what '
        'that is or why they need to.\n\n'
        'Priority fixes:\n'
        '1. Add a 2-3 sentence explainer above the list: "A federation is a group that manages your wallet '
        'and keeps your money safe. Think of it like choosing your bank."\n'
        '2. Make federation cards tappable with a detail view: description, guardian info, member count, '
        'terms summary.\n'
        '3. Fix text truncation (technical, easy).\n\n'
        'Heuristics: H2 (match with real world), H5 (error prevention)')

    finding_slide(
        prs, 3, 'Joining a Federation -- Errors, Loops & a Disappearing Exit',
        'They tap Join. A Terms of Service screen appears, titled simply Welcome, with ToS as an external PDF.',
        SCREENSHOTS + '03-join-tos.png',
        criticals=[
            'Screen title is "Welcome" -- no signal this is a legal consent moment. ToS not shown in-app.',
            'Error banner is generic, has no retry button, and persists across multiple screens.',
            '"View public federations" on QR scanner loops back to the same list.',
            '"Maybe Later" disappears mid-flow. Back button exits the app entirely.',
        ],
        moderates=[
            "Camera permission screen doesn\'t clarify whether Continue grants or defers permission.",
            'QR scanner has no instruction text -- what should the user point their camera at?',
        ],
        positives=[],
    )
    set_notes(list(prs.slides)[-1],
        'This section contains the most dangerous UX failures -- the ones that can strand a user with '
        'no path forward:\n\n'
        '- The ToS moment: label it clearly ("Before you join"), show key terms in-app (scrollable), '
        'link to full PDF as secondary action.\n'
        '- Error banner: add a Retry button. Dismiss automatically on success. Don\'t let it persist '
        'across screens.\n'
        '- The "walled in" problem: if joining fails, users must always be able to get back to the home '
        'screen or the federation list. Back button should never exit the app from within a flow.\n'
        '- QR scanner: one line of instructional copy ("Scan a federation QR code from an admin or '
        'the Fedi website") would eliminate the confusion entirely.\n\n'
        'Heuristics: H1 (system status), H3 (user control), H5 (error prevention), H9 (error recovery)')

    finding_slide(
        prs, 4, 'Arriving in the App',
        'If they make it through -- they land on home with a modal already waiting.',
        SCREENSHOTS + '05-home-username-modal.png',
        criticals=[
            'Display name auto-assigned as "ambitious wolf" without asking. Jarring for a chat + payments app.',
        ],
        moderates=[
            '"ambitious wolf #267f" -- the hash suffix is technical noise with no explanation.',
            'Home screen mixes three content types at equal visual weight.',
            'Fedi Global banner has an unlabelled QR icon -- not obviously tappable.',
        ],
        positives=[],
    )
    set_notes(list(prs.slides)[-1],
        'The auto-assigned display name is a significant trust and identity issue for a payments app. '
        'Users are about to receive and send real money -- having an embarrassing or random display name '
        'undermines confidence.\n\n'
        'Fix: prompt for a display name during onboarding, before entering the app. Even a simple '
        '"What should people call you?" screen.\n\n'
        'Home screen hierarchy: the three content zones (federation banner, Community News, Mini Apps) '
        'need visual separation -- section headers, spacing, or background differentiation.\n\n'
        'Hash suffix (#267f): this is a Nostr public key fragment. It should not be surfaced to users '
        'without an explanation, or should be hidden entirely behind a "technical details" toggle.\n\n'
        'Heuristics: H2 (match with real world), H8 (aesthetic and minimalist design)')

    slide_nav(prs)
    set_notes(list(prs.slides)[-1],
        'The unlabelled bottom navigation is a systemic issue, not a cosmetic one. '
        'Android Material Design guidelines recommend labels on bottom nav items precisely because '
        'icons alone require learned knowledge. For a new user, four mystery icons mean four things '
        'to trial-and-error.\n\n'
        'Promotional bottom sheets on every tab open are particularly damaging: the first thing a '
        'user sees on switching to any section is a sales message they must dismiss. This pattern '
        'trains users to dismiss quickly without reading, which is dangerous if a future sheet '
        'contains important information.\n\n'
        'The Wallet tab redirect is a critical wayfinding failure: if you cannot use the wallet '
        'without a federation, the tab should be disabled with an explanation, not silently redirect.\n\n'
        'Heuristics: H1, H3, H4, H8')

    slide_settings(prs)
    set_notes(list(prs.slides)[-1],
        'Settings is structurally sound -- the content is the right content. The problems are '
        'presentation issues:\n\n'
        '- Unlabelled QR code at the top: this is almost certainly a wallet receive address or a '
        'profile link. Label it. "Your receive QR code" or "Share your profile" -- one line.\n'
        '- "#267f" hash suffix: hide or explain (see slide 4 notes).\n'
        '- "Nostr Details": this term is meaningful to a small subset of users. For everyone else, '
        'rename to "Social identity" or tuck under an Advanced section.\n'
        '- "Share logs": this is a support/debugging feature. It should look tappable (chevron, '
        'icon) and sit in a clearly marked Support section.\n\n'
        'Positives: the settings content itself is well-chosen and well-organised. The bones are good.')

    slide_priority_matrix(prs)
    set_notes(list(prs.slides)[-1],
        'Reading the matrix:\n'
        '- Filled circle = critical: these should be in the next sprint.\n'
        '- Triangle = moderate: address in the following 1-2 sprints.\n'
        '- Check = works well: preserve these patterns in future design decisions.\n\n'
        'Top 3 critical fixes by effort-to-impact ratio:\n'
        '1. Add a federation explainer before the discovery screen (copy change, low effort, high impact)\n'
        '2. Add labels to bottom navigation (design change, low-medium effort, high impact)\n'
        '3. Fix the "walled in" back-button behaviour (engineering fix, medium effort, critical safety issue)')

    slide_legend(prs)
    set_notes(list(prs.slides)[-1],
        'Nielsen Norman Group 10 Usability Heuristics (1994, updated 2020) are the industry standard '
        'framework for expert UX evaluation. Each finding is tagged to the heuristic it violates, '
        'providing a principled basis for prioritisation rather than subjective preference.\n\n'
        'Full reference: nngroup.com/articles/ten-usability-heuristics/')

    section_divider(prs, 'Assumptions + Research Questions',
                    "What we\'re betting on -- and how to test each bet")
    set_notes(list(prs.slides)[-1],
        'Transition note: the UX audit tells us what is broken today. The assumptions section '
        'tells us what we need to validate before we know how to fix it correctly for the target audience.\n\n'
        '16 assumptions are mapped across 4 categories: mental model/comprehension, trust & safety, '
        'onboarding & flows, and strategy & governance. Each has an explicit risk level, evidence '
        'status, and research method recommendation in the Assumption Register.')

    slide_assumptions_overview(prs)
    set_notes(list(prs.slides)[-1],
        'Overview of all 16 assumptions. Colour coding:\n'
        '- Blue: mental model / comprehension assumptions\n'
        '- Red: trust & safety assumptions\n'
        '- Green: onboarding & flows assumptions\n'
        '- Yellow: strategy & governance assumptions\n\n'
        'All 16 currently have no or weak evidence. The assumption register (later in this deck) '
        'prioritises them by risk and maps each to a specific research method and timebox.')

    slides_rq_detail(prs)
    # No per-slide notes here -- detail slides are self-contained

    section_divider(prs, 'Assumption Register',
                    'Prioritised research backlog: risk, evidence status, method & timebox')
    set_notes(list(prs.slides)[-1],
        'The Assumption Register is the live artefact that should be maintained in Notion or Sheets '
        'after this conversation. It answers the question "what are we currently betting on and '
        'what is the plan to validate it?" for any stakeholder at any time.\n\n'
        'Risk levels:\n'
        '- H (High): getting this wrong has significant product or user safety consequences\n'
        '- M (Medium): getting this wrong creates friction or missed retention\n'
        '- L (Low): relatively safe to ship and iterate on\n\n'
        'Evidence status will evolve from None/Weak toward Moderate/Strong as studies run.')

    slides_assumption_register(prs)

    section_divider(prs, '30 / 60 / 90 Day Plan',
                    'Solo researcher + fast shipping environment')
    set_notes(list(prs.slides)[-1],
        'This plan is designed for a researcher who is:\n'
        '- The only UX researcher in the org\n'
        '- Working in a fast-shipping, resource-constrained environment\n'
        '- Responsible for covering multiple regions and a complex product\n\n'
        'The phased approach prioritises: credibility first, then coverage, then compounding value. '
        'It mirrors the structure Anwar has successfully executed at previous orgs.')

    slide_9090_summary(prs)
    slide_30(prs)
    slide_60(prs)
    slide_90(prs)

    out = '/Users/annie/fedi-ux-audit-deck/fedi-ux-audit.pptx'
    prs.save(out)
    print('Saved:', out)
    print('Total slides:', len(prs.slides))

if __name__ == '__main__':
    build()
